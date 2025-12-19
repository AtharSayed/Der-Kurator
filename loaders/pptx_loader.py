# loaders/pptx_loader.py
from pptx import Presentation
from loaders.base_loader import BaseLoader

class PptxLoader(BaseLoader):
    def load(self, file_path):
        prs = Presentation(file_path)
        documents = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

            slide_text = "\n".join(texts).strip()
            if slide_text:
                documents.append({
                    "content": slide_text,
                    "metadata": {
                        "source": file_path,
                        "type": "pptx",
                        "unit": "slide",
                        "index": i + 1
                    }
                })
        return documents